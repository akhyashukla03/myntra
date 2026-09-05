import sys
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# Initialize Presentation with 16:9 Widescreen dimensions
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

blank_slide_layout = prs.slide_layouts[6]

# Authentic Myntra Corporate Brand Palette Definitions
COLOR_WHITE = RGBColor(255, 255, 255)
COLOR_OFFWHITE = RGBColor(245, 245, 246)          # Official Myntra App Background #F5F5F6
COLOR_BG_CARD = RGBColor(255, 255, 255)           # Pure White Card #FFFFFF
COLOR_BORDER = RGBColor(234, 234, 236)            # Official Myntra Card Border #EAEAEC
COLOR_TEXT_PRIMARY = RGBColor(40, 44, 63)         # Official Myntra Deep Charcoal #282C3F
COLOR_TEXT_MUTED = RGBColor(83, 87, 102)          # Official Myntra Subtext Slate #535766
COLOR_MYNTRA_PINK = RGBColor(255, 63, 108)        # Signature Myntra Pink #FF3F6C
COLOR_PURPLE_PILL = RGBColor(255, 63, 108)        # Myntra Pink Primary Header Pill
COLOR_SUBTITLE = RGBColor(255, 63, 108)           # Myntra Pink Subtitle Accent
COLOR_SYNTHESIS_BG = RGBColor(255, 240, 244)      # Myntra Soft Pink Glow #FFF0F4
COLOR_SYNTHESIS_BORDER = RGBColor(255, 194, 209)  # Myntra Soft Pink Border #FFC2D1
COLOR_SYNTHESIS_TEXT = RGBColor(40, 44, 63)       # Myntra Deep Charcoal #282C3F
COLOR_PHONE_BG = RGBColor(40, 44, 63)             # Myntra Deep Charcoal #282C3F
COLOR_MYNTRA_GREEN = RGBColor(3, 166, 133)        # Myntra Verified Trust Green #03A685
COLOR_LIGHT_GRAY = RGBColor(248, 249, 250)

# Strict Font Size Standards (Adhering to strict >= 14pt Presentation Guidelines)
FONT_BANNER = Pt(14)
FONT_TITLE = Pt(18)
FONT_SUBTITLE = Pt(14)
FONT_CARD_HEADER = Pt(15)
FONT_BODY = Pt(14)
FONT_METRIC_VAL = Pt(24)
FONT_METRIC_LBL = Pt(14)

def add_card_shape(slide, left, top, width, height, bg_rgb, border_rgb=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_rgb
    if border_rgb:
        shape.line.color.rgb = border_rgb
        shape.line.width = Pt(1.5)
    else:
        shape.line.fill.background()
    return shape

def set_bullet_indent(p, left_margin_pt=18, hanging_indent_pt=14):
    """
    Sets true OpenXML paragraph left margin (marL) and negative first-line indent (indent)
    so that multi-line wrapped bullet text aligns cleanly under the text block.
    """
    pPr = p._p.get_or_add_pPr()
    pPr.set('marL', str(int(Pt(left_margin_pt))))
    pPr.set('indent', str(int(-Pt(hanging_indent_pt))))

# Import slide data from python logic
from data_pptx import SLIDES_DATA

for data in SLIDES_DATA:
    slide = prs.slides.add_slide(blank_slide_layout)
    s_num = data["slideNumber"]
    
    # 1. Top Category Pill Banner (14pt Bold)
    banner_width = Inches(7.5)
    banner_height = Inches(0.42)
    banner_left = Inches((13.333 - 7.5) / 2)
    banner_top = Inches(0.18)
    
    pill_shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, banner_left, banner_top, banner_width, banner_height)
    pill_shape.fill.solid()
    pill_shape.fill.fore_color.rgb = COLOR_PURPLE_PILL
    pill_shape.line.fill.background()
    
    tf = pill_shape.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Inches(0.1)
    tf.margin_right = Inches(0.1)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    
    p = tf.paragraphs[0]
    p.text = data["topBanner"]
    p.alignment = PP_ALIGN.CENTER
    p.font.size = FONT_BANNER
    p.font.bold = True
    p.font.color.rgb = COLOR_WHITE
    p.space_before = Pt(0)
    p.space_after = Pt(0)

    # 2. Main Header & Brand Title Block (18pt Title, 14pt Subtitle)
    title_box = slide.shapes.add_textbox(Inches(0.65), Inches(0.65), Inches(10.8), Inches(0.85))
    tf = title_box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0)
    tf.margin_right = Inches(0)
    tf.margin_top = Inches(0)
    tf.margin_bottom = Inches(0)
    
    p = tf.paragraphs[0]
    p.text = data["title"]
    p.font.size = FONT_TITLE
    p.font.bold = True
    p.font.color.rgb = COLOR_TEXT_PRIMARY
    p.space_after = Pt(2)
    
    p2 = tf.add_paragraph()
    p2.text = data["subtitle"]
    p2.font.size = FONT_SUBTITLE
    p2.font.bold = True
    p2.font.color.rgb = COLOR_SUBTITLE
    p2.space_before = Pt(2)
    p2.space_after = Pt(4)
    
    brand_box = slide.shapes.add_textbox(Inches(11.4), Inches(0.65), Inches(1.3), Inches(0.5))
    tf_b = brand_box.text_frame
    tf_b.margin_left = Inches(0)
    tf_b.margin_right = Inches(0)
    tf_b.margin_top = Inches(0)
    tf_b.margin_bottom = Inches(0)
    p_b = tf_b.paragraphs[0]
    p_b.text = "myntra"
    p_b.font.size = Pt(24)
    p_b.font.bold = True
    p_b.font.color.rgb = COLOR_MYNTRA_PINK
    p_b.alignment = PP_ALIGN.RIGHT

    # 3. Slide Content Area Layouts (Full 5.68" vertical canvas from top 1.55")
    card_top = Inches(1.55)
    card_height = Inches(5.68)
    
    # ==========================================
    # SLIDE 1: STRATEGIC GOAL & 5-STEP FINANCIAL DERIVATION WATERFALL
    # ==========================================
    if s_num == 1:
        # Left Card: Strategic Goal & Core Problem (Top area)
        left_w = Inches(4.3)
        top_h = Inches(4.15)
        add_card_shape(slide, Inches(0.65), card_top, left_w, top_h, COLOR_BG_CARD, COLOR_BORDER)
        box = slide.shapes.add_textbox(Inches(0.77), card_top + Inches(0.08), left_w - Inches(0.24), top_h - Inches(0.16))
        tf = box.text_frame; tf.word_wrap = True
        tf.margin_left = Inches(0.06); tf.margin_right = Inches(0.06)
        
        p = tf.paragraphs[0]
        p.text = "🎯 " + data["leftCard"]["title"]
        p.font.size = FONT_CARD_HEADER
        p.font.bold = True
        p.font.color.rgb = COLOR_PURPLE_PILL
        p.space_after = Pt(4)
        for bold_text, reg_text in data["leftCard"]["bullets"]:
            p = tf.add_paragraph()
            p.font.size = FONT_BODY
            p.line_spacing = 1.15
            p.space_before = Pt(2); p.space_after = Pt(3)
            set_bullet_indent(p, left_margin_pt=16, hanging_indent_pt=12)
            r0 = p.add_run(); r0.text = "• "; r0.font.color.rgb = COLOR_MYNTRA_PINK; r0.font.bold = True
            r1 = p.add_run(); r1.text = bold_text + " "; r1.font.bold = True; r1.font.color.rgb = COLOR_TEXT_PRIMARY
            r2 = p.add_run(); r2.text = reg_text; r2.font.color.rgb = COLOR_TEXT_PRIMARY

        # Mid Card: 5-Step Financial Derivation Waterfall
        mid_left = Inches(5.1)
        mid_w = Inches(4.8)
        add_card_shape(slide, mid_left, card_top, mid_w, top_h, COLOR_SYNTHESIS_BG, COLOR_SYNTHESIS_BORDER)
        box_mid = slide.shapes.add_textbox(mid_left + Inches(0.12), card_top + Inches(0.08), mid_w - Inches(0.24), top_h - Inches(0.16))
        tf_mid = box_mid.text_frame; tf_mid.word_wrap = True
        tf_mid.margin_left = Inches(0.06); tf_mid.margin_right = Inches(0.06)
        
        p = tf_mid.paragraphs[0]
        p.text = "💰 " + data["midCard"]["title"]
        p.font.size = FONT_CARD_HEADER
        p.font.bold = True
        p.font.color.rgb = COLOR_PURPLE_PILL
        p.space_after = Pt(4)
        
        for bold_text, reg_text in data["midCard"]["bullets"]:
            p_st = tf_mid.add_paragraph()
            p_st.font.size = FONT_BODY
            p_st.line_spacing = 1.15
            p_st.space_before = Pt(2); p_st.space_after = Pt(2)
            set_bullet_indent(p_st, left_margin_pt=14, hanging_indent_pt=0)
            r1 = p_st.add_run(); r1.text = bold_text + " "; r1.font.bold = True; r1.font.color.rgb = COLOR_TEXT_PRIMARY
            r2 = p_st.add_run(); r2.text = reg_text; r2.font.color.rgb = COLOR_TEXT_PRIMARY

        # Right Phone Screen Mockup
        phone_left = Inches(10.05)
        phone_width = Inches(2.63)
        add_card_shape(slide, phone_left, card_top, phone_width, top_h, COLOR_PHONE_BG, RGBColor(15, 23, 42))
        figma_img = data.get("figmaImage")
        if figma_img and os.path.exists(figma_img):
            slide.shapes.add_picture(figma_img, phone_left + Inches(0.06), card_top + Inches(0.06), phone_width - Inches(0.12), top_h - Inches(0.12))

        # Bottom: 4 Macro Metrics Boxes spanning the full width
        macro_top = card_top + top_h + Inches(0.15)
        macro_w = Inches(2.88)
        macro_gap = Inches(0.17)
        macro_h = Inches(1.35)
        for m_idx, (m_val, m_lbl) in enumerate(data["macroMetrics"]):
            m_left = Inches(0.65) + m_idx * (macro_w + macro_gap)
            add_card_shape(slide, m_left, macro_top, macro_w, macro_h, COLOR_WHITE, COLOR_BORDER)
            m_box = slide.shapes.add_textbox(m_left, macro_top + Inches(0.08), macro_w, macro_h - Inches(0.16))
            tf_m = m_box.text_frame; tf_m.word_wrap = True
            p_val = tf_m.paragraphs[0]
            p_val.text = m_val
            p_val.alignment = PP_ALIGN.CENTER
            p_val.font.size = FONT_METRIC_VAL
            p_val.font.bold = True
            p_val.font.color.rgb = COLOR_MYNTRA_PINK
            p_lbl = tf_m.add_paragraph()
            p_lbl.text = m_lbl
            p_lbl.alignment = PP_ALIGN.CENTER
            p_lbl.font.size = FONT_METRIC_LBL
            p_lbl.font.bold = True
            p_lbl.font.color.rgb = COLOR_TEXT_MUTED
            p_lbl.space_before = Pt(2)

    # ==========================================
    # SLIDE 2: HYPOTHESIS, FUNNEL & COMPETITOR TEARDOWN
    # ==========================================
    elif s_num == 2:
        # Top Left: Hypothesis Card
        add_card_shape(slide, Inches(0.65), card_top, Inches(4.5), Inches(1.6), COLOR_WHITE, COLOR_BORDER)
        box = slide.shapes.add_textbox(Inches(0.77), card_top + Inches(0.08), Inches(4.26), Inches(1.44))
        tf = box.text_frame; tf.word_wrap = True
        tf.margin_left = Inches(0.06); tf.margin_right = Inches(0.06)
        p = tf.paragraphs[0]
        p.text = "💡 THE CORE HYPOTHESIS"
        p.font.size = FONT_CARD_HEADER; p.font.bold = True; p.font.color.rgb = COLOR_PURPLE_PILL; p.space_after = Pt(2)
        p_body = tf.add_paragraph()
        p_body.text = data["hypothesisBox"]
        p_body.font.size = FONT_BODY; p_body.line_spacing = 1.15; p_body.font.color.rgb = COLOR_TEXT_PRIMARY

        # Top Right: Current Discovery Funnel (4 Process Steps)
        flow_left = Inches(5.3)
        flow_w = Inches(7.383)
        add_card_shape(slide, flow_left, card_top, flow_w, Inches(1.6), COLOR_WHITE, COLOR_BORDER)
        box_hdr = slide.shapes.add_textbox(flow_left + Inches(0.12), card_top + Inches(0.06), flow_w - Inches(0.24), Inches(0.35))
        tf_h = box_hdr.text_frame; tf_h.word_wrap = True
        p_h = tf_h.paragraphs[0]
        p_h.text = "🔄 CURRENT DISCOVERY FLOW (THE GRAVEYARD LOOP)"
        p_h.font.size = FONT_CARD_HEADER; p_h.font.bold = True; p_h.font.color.rgb = COLOR_PURPLE_PILL

        step_w = Inches(1.72)
        step_gap = Inches(0.1)
        for s_idx, (step_title, step_desc) in enumerate(data["discoveryFunnel"]):
            s_left = flow_left + Inches(0.12) + s_idx * (step_w + step_gap)
            add_card_shape(slide, s_left, card_top + Inches(0.42), step_w, Inches(1.08), COLOR_SYNTHESIS_BG, COLOR_SYNTHESIS_BORDER)
            s_box = slide.shapes.add_textbox(s_left + Inches(0.04), card_top + Inches(0.44), step_w - Inches(0.08), Inches(1.02))
            tf_s = s_box.text_frame; tf_s.word_wrap = True
            tf_s.margin_left = Inches(0.02); tf_s.margin_right = Inches(0.02)
            p_st = tf_s.paragraphs[0]
            p_st.text = step_title
            p_st.font.size = FONT_BODY; p_st.font.bold = True; p_st.font.color.rgb = COLOR_MYNTRA_PINK; p_st.space_after = Pt(1)
            p_sd = tf_s.add_paragraph()
            p_sd.text = step_desc
            p_sd.font.size = FONT_BODY; p_sd.line_spacing = 1.1; p_sd.font.color.rgb = COLOR_TEXT_PRIMARY

        # Bottom Left: Competitor Whitespace Teardown
        comp_top = card_top + Inches(1.75)
        comp_w = Inches(5.8)
        comp_h = card_height - Inches(1.75)
        add_card_shape(slide, Inches(0.65), comp_top, comp_w, comp_h, COLOR_BG_CARD, COLOR_BORDER)
        box_comp = slide.shapes.add_textbox(Inches(0.77), comp_top + Inches(0.08), comp_w - Inches(0.24), comp_h - Inches(0.16))
        tf_c = box_comp.text_frame; tf_c.word_wrap = True
        tf_c.margin_left = Inches(0.06); tf_c.margin_right = Inches(0.06)
        p = tf_c.paragraphs[0]
        p.text = "🌐 GLOBAL & DOMESTIC COMPETITOR TEARDOWN"
        p.font.size = FONT_CARD_HEADER; p.font.bold = True; p.font.color.rgb = COLOR_PURPLE_PILL; p.space_after = Pt(4)

        for comp_platform, comp_text in data["competitorTeardown"]:
            p_cp = tf_c.add_paragraph()
            p_cp.font.size = FONT_BODY; p_cp.line_spacing = 1.15
            p_cp.space_before = Pt(2); p_cp.space_after = Pt(2)
            set_bullet_indent(p_cp, left_margin_pt=14, hanging_indent_pt=0)
            r1 = p_cp.add_run(); r1.text = "• " + comp_platform + " "; r1.font.bold = True; r1.font.color.rgb = COLOR_MYNTRA_PINK if "Myntra" in comp_platform else COLOR_TEXT_PRIMARY
            r2 = p_cp.add_run(); r2.text = comp_text; r2.font.color.rgb = COLOR_TEXT_PRIMARY

        # Bottom Right: 3 Friction Cards
        f_top = comp_top
        f_card_w = Inches(6.083)
        f_card_h = (comp_h - Inches(0.2)) / 3.0
        for f_idx, (f_title, f_desc) in enumerate(data["frictionCards"]):
            f_cur_top = f_top + f_idx * (f_card_h + Inches(0.1))
            add_card_shape(slide, Inches(6.6), f_cur_top, f_card_w, f_card_h, COLOR_SYNTHESIS_BG, COLOR_SYNTHESIS_BORDER)
            f_box = slide.shapes.add_textbox(Inches(6.72), f_cur_top + Inches(0.06), f_card_w - Inches(0.24), f_card_h - Inches(0.12))
            tf_f = f_box.text_frame; tf_f.word_wrap = True
            tf_f.margin_left = Inches(0.06); tf_f.margin_right = Inches(0.06)
            p_ft = tf_f.paragraphs[0]
            p_ft.text = "⚠️ " + f_title
            p_ft.font.size = FONT_CARD_HEADER; p_ft.font.bold = True; p_ft.font.color.rgb = COLOR_MYNTRA_PINK; p_ft.space_after = Pt(2)
            p_fb = tf_f.add_paragraph()
            p_fb.text = f_desc
            p_fb.font.size = FONT_BODY; p_fb.line_spacing = 1.12; p_fb.font.color.rgb = COLOR_TEXT_PRIMARY

    # ==========================================
    # SLIDE 3: THINKING EVOLUTION & NLP DISCOVERY ENGINE
    # ==========================================
    elif s_num == 3:
        # Top Card: 2-Sentence Thinking Evolution Narrative
        evol_h = Inches(1.4)
        add_card_shape(slide, Inches(0.65), card_top, Inches(12.033), evol_h, COLOR_SYNTHESIS_BG, COLOR_SYNTHESIS_BORDER)
        box_ev = slide.shapes.add_textbox(Inches(0.77), card_top + Inches(0.08), Inches(11.793), evol_h - Inches(0.16))
        tf_ev = box_ev.text_frame; tf_ev.word_wrap = True
        tf_ev.margin_left = Inches(0.06); tf_ev.margin_right = Inches(0.06)
        p_eh = tf_ev.paragraphs[0]
        p_eh.text = "🔄 STRATEGIC THINKING EVOLUTION & RESEARCH DISCOVERY"
        p_eh.font.size = FONT_CARD_HEADER; p_eh.font.bold = True; p_eh.font.color.rgb = COLOR_PURPLE_PILL; p_eh.space_after = Pt(2)
        
        p_eb = tf_ev.add_paragraph()
        p_eb.font.size = FONT_BODY; p_eb.line_spacing = 1.15
        set_bullet_indent(p_eb, left_margin_pt=14, hanging_indent_pt=0)
        r1 = p_eb.add_run(); r1.text = data["evolutionNarrative"]["initialHypothesis"] + "  |  "; r1.font.color.rgb = COLOR_TEXT_PRIMARY
        r2 = p_eb.add_run(); r2.text = data["evolutionNarrative"]["dataFinding"] + "  |  "; r2.font.bold = True; r2.font.color.rgb = COLOR_MYNTRA_PINK
        r3 = p_eb.add_run(); r3.text = data["evolutionNarrative"]["strategicPivot"]; r3.font.bold = True; r3.font.color.rgb = COLOR_MYNTRA_GREEN

        # Mid: 3 NLP Prompts & Findings side-by-side cards
        nlp_top = card_top + evol_h + Inches(0.12)
        nlp_w = Inches(3.9)
        nlp_gap = Inches(0.16)
        nlp_h = Inches(2.55)
        for p_idx, (q_txt, ans_txt) in enumerate(data["nlpPrompts"]):
            p_left = Inches(0.65) + p_idx * (nlp_w + nlp_gap)
            add_card_shape(slide, p_left, nlp_top, nlp_w, nlp_h, COLOR_BG_CARD, COLOR_BORDER)
            box_p = slide.shapes.add_textbox(p_left + Inches(0.08), nlp_top + Inches(0.06), nlp_w - Inches(0.16), nlp_h - Inches(0.12))
            tf_p = box_p.text_frame; tf_p.word_wrap = True
            tf_p.margin_left = Inches(0.06); tf_p.margin_right = Inches(0.06)
            p = tf_p.paragraphs[0]
            p.text = "💬 " + q_txt
            p.font.size = FONT_CARD_HEADER; p.font.bold = True; p.font.color.rgb = COLOR_MYNTRA_PINK; p.space_after = Pt(3)
            p_s = tf_p.add_paragraph()
            p_s.text = ans_txt
            p_s.font.size = FONT_BODY; p_s.line_spacing = 1.14; p_s.font.color.rgb = COLOR_TEXT_PRIMARY

        # Bottom: 3-Step NLP Pipeline & Testable Link
        pipe_top = nlp_top + nlp_h + Inches(0.12)
        pipe_h = card_height - (evol_h + nlp_h + Inches(0.24))
        add_card_shape(slide, Inches(0.65), pipe_top, Inches(12.033), pipe_h, COLOR_WHITE, COLOR_BORDER)
        box_pipe = slide.shapes.add_textbox(Inches(0.77), pipe_top + Inches(0.06), Inches(11.793), pipe_h - Inches(0.12))
        tf_pipe = box_pipe.text_frame; tf_pipe.word_wrap = True
        tf_pipe.margin_left = Inches(0.06); tf_pipe.margin_right = Inches(0.06)
        p = tf_pipe.paragraphs[0]
        p.text = "🔬 3-STEP NLP DISCOVERY ENGINE PIPELINE & TESTABLE WORKFLOW"
        p.font.size = FONT_CARD_HEADER; p.font.bold = True; p.font.color.rgb = COLOR_PURPLE_PILL; p.space_after = Pt(2)

        p_stages = tf_pipe.add_paragraph()
        p_stages.font.size = FONT_BODY; p_stages.line_spacing = 1.12
        for s_idx, (s_title, s_desc) in enumerate(data["corpusFunnel"]):
            r0 = p_stages.add_run(); r0.text = f"[{s_title}]: "; r0.font.bold = True; r0.font.color.rgb = COLOR_MYNTRA_PINK
            r1 = p_stages.add_run(); r1.text = s_desc + ("   |   " if s_idx < len(data["corpusFunnel"])-1 else "")
            r1.font.color.rgb = COLOR_TEXT_PRIMARY

        p_link = tf_pipe.add_paragraph()
        p_link.font.size = FONT_BODY; p_link.space_before = Pt(2)
        r_l1 = p_link.add_run(); r_l1.text = "🔗 Live Testable Workflow: "; r_l1.font.bold = True; r_l1.font.color.rgb = COLOR_MYNTRA_GREEN
        r_l2 = p_link.add_run(); r_l2.text = data["discoveryEngineUrl"]; r_l2.font.bold = True; r_l2.font.color.rgb = COLOR_TEXT_PRIMARY

    # ==========================================
    # SLIDE 4: 4-QUADRANT TARGET SEGMENT CANVAS
    # ==========================================
    elif s_num == 4:
        quad_w = Inches(5.9)
        quad_h = (card_height - Inches(0.18)) / 2.0
        positions = [
            (Inches(0.65), card_top),
            (Inches(6.78), card_top),
            (Inches(0.65), card_top + quad_h + Inches(0.18)),
            (Inches(6.78), card_top + quad_h + Inches(0.18))
        ]
        for q_idx, q_data in enumerate(data["quadrants"]):
            q_left, q_top = positions[q_idx]
            add_card_shape(slide, q_left, q_top, quad_w, quad_h, COLOR_BG_CARD, COLOR_BORDER)
            box = slide.shapes.add_textbox(q_left + Inches(0.12), q_top + Inches(0.08), quad_w - Inches(0.24), quad_h - Inches(0.16))
            tf = box.text_frame; tf.word_wrap = True
            tf.margin_left = Inches(0.06); tf.margin_right = Inches(0.06)
            p = tf.paragraphs[0]
            p.text = f"{q_data['icon']} {q_data['title']}"
            p.font.size = FONT_CARD_HEADER; p.font.bold = True; p.font.color.rgb = COLOR_PURPLE_PILL; p.space_after = Pt(3)
            for b_lbl, b_txt in q_data["bullets"]:
                p = tf.add_paragraph()
                p.font.size = FONT_BODY; p.line_spacing = 1.15
                p.space_before = Pt(1); p.space_after = Pt(2)
                set_bullet_indent(p, left_margin_pt=14, hanging_indent_pt=10)
                r1 = p.add_run(); r1.text = "• " + b_lbl + " "; r1.font.bold = True; r1.font.color.rgb = COLOR_TEXT_PRIMARY
                r2 = p.add_run(); r2.text = b_txt; r2.font.color.rgb = COLOR_TEXT_PRIMARY

    # ==========================================
    # SLIDE 5: QUALITATIVE USER RESEARCH (2x3 GRID - 6 CARDS)
    # ==========================================
    elif s_num == 5:
        # Top Methodology Banner
        meth_h = Inches(0.48)
        add_card_shape(slide, Inches(0.65), card_top, Inches(12.033), meth_h, COLOR_SYNTHESIS_BG, COLOR_SYNTHESIS_BORDER)
        box_m = slide.shapes.add_textbox(Inches(0.77), card_top + Inches(0.04), Inches(11.793), meth_h - Inches(0.08))
        tf_m = box_m.text_frame; tf_m.word_wrap = True
        p_m = tf_m.paragraphs[0]
        p_m.text = "📋 " + data.get("researchMethodology", "Qualitative User Interviews (N=9 in-depth sessions across Tier 1 & Tier 2 shoppers)")
        p_m.font.size = FONT_BANNER; p_m.font.bold = True; p_m.font.color.rgb = COLOR_PURPLE_PILL

        # 2x3 Grid of 6 User Cards
        grid_top = card_top + meth_h + Inches(0.12)
        grid_w = Inches(3.9)
        grid_gap_x = Inches(0.16)
        grid_h = (card_height - meth_h - Inches(0.24)) / 2.0
        grid_gap_y = Inches(0.12)

        for u_idx, user in enumerate(data["userCards"]):
            col = u_idx % 3
            row = u_idx // 3
            u_left = Inches(0.65) + col * (grid_w + grid_gap_x)
            u_top = grid_top + row * (grid_h + grid_gap_y)
            
            add_card_shape(slide, u_left, u_top, grid_w, grid_h, COLOR_BG_CARD, COLOR_BORDER)
            box_u = slide.shapes.add_textbox(u_left + Inches(0.08), u_top + Inches(0.06), grid_w - Inches(0.16), grid_h - Inches(0.12))
            tf_u = box_u.text_frame; tf_u.word_wrap = True
            tf_u.margin_left = Inches(0.06); tf_u.margin_right = Inches(0.06)
            
            # Participant Header
            p_uh = tf_u.paragraphs[0]
            p_uh.text = "👤 " + user["id"] + " (" + user["demographics"] + ")"
            p_uh.font.size = FONT_CARD_HEADER; p_uh.font.bold = True; p_uh.font.color.rgb = COLOR_PURPLE_PILL; p_uh.space_after = Pt(2)
            
            # Authentic Quote Box
            p_q = tf_u.add_paragraph()
            p_q.font.size = FONT_BODY; p_q.line_spacing = 1.12; p_q.space_after = Pt(2)
            r_q1 = p_q.add_run(); r_q1.text = 'Quote: '; r_q1.font.bold = True; r_q1.font.color.rgb = COLOR_MYNTRA_PINK
            r_q2 = p_q.add_run(); r_q2.text = user["quote"]; r_q2.font.italic = True; r_q2.font.color.rgb = COLOR_TEXT_PRIMARY

            # Key PM Insight
            p_ins = tf_u.add_paragraph()
            p_ins.font.size = FONT_BODY; p_ins.line_spacing = 1.12; p_ins.space_after = Pt(0)
            r_i1 = p_ins.add_run(); r_i1.text = user["insight"]; r_i1.font.bold = True; r_i1.font.color.rgb = COLOR_MYNTRA_GREEN

    # ==========================================
    # SLIDE 6: PROBLEM FRAMING (5 GOLD-STANDARD PM QUESTIONS)
    # ==========================================
    elif s_num == 6:
        # Top Row: 3 Questions (Q1, Q2, Q3)
        top_h = (card_height - Inches(0.18)) / 2.0
        q_w = Inches(3.9)
        q_gap = Inches(0.16)
        for q_idx in range(3):
            q_data = data["pmQuestions"][q_idx]
            q_left = Inches(0.65) + q_idx * (q_w + q_gap)
            add_card_shape(slide, q_left, card_top, q_w, top_h, COLOR_BG_CARD, COLOR_BORDER)
            box = slide.shapes.add_textbox(q_left + Inches(0.08), card_top + Inches(0.06), q_w - Inches(0.16), top_h - Inches(0.12))
            tf = box.text_frame; tf.word_wrap = True
            tf.margin_left = Inches(0.06); tf.margin_right = Inches(0.06)
            p = tf.paragraphs[0]
            p.text = f"Q{q_idx+1}: {q_data['q']}"
            p.font.size = FONT_CARD_HEADER; p.font.bold = True; p.font.color.rgb = COLOR_PURPLE_PILL; p.space_after = Pt(3)
            
            p_ans = tf.add_paragraph()
            p_ans.text = q_data["ans"]
            p_ans.font.size = FONT_BODY; p_ans.line_spacing = 1.14; p_ans.font.color.rgb = COLOR_TEXT_PRIMARY

        # Bottom Row: 2 Questions (Q4, Q5)
        bot_top = card_top + top_h + Inches(0.18)
        bot_w = Inches(5.9)
        bot_gap = Inches(0.23)
        for q_idx in range(3, 5):
            q_data = data["pmQuestions"][q_idx]
            q_left = Inches(0.65) + (q_idx - 3) * (bot_w + bot_gap)
            add_card_shape(slide, q_left, bot_top, bot_w, top_h, COLOR_BG_CARD, COLOR_BORDER)
            box = slide.shapes.add_textbox(q_left + Inches(0.08), bot_top + Inches(0.06), bot_w - Inches(0.16), top_h - Inches(0.12))
            tf = box.text_frame; tf.word_wrap = True
            tf.margin_left = Inches(0.06); tf.margin_right = Inches(0.06)
            p = tf.paragraphs[0]
            p.text = f"Q{q_idx+1}: {q_data['q']}"
            p.font.size = FONT_CARD_HEADER; p.font.bold = True; p.font.color.rgb = COLOR_PURPLE_PILL; p.space_after = Pt(3)
            
            if "customerValue" in q_data:
                p_c = tf.add_paragraph()
                p_c.font.size = FONT_BODY; p_c.line_spacing = 1.14; p_c.space_after = Pt(2)
                r_c1 = p_c.add_run(); r_c1.text = "Customer Value: "; r_c1.font.bold = True; r_c1.font.color.rgb = COLOR_TEXT_PRIMARY
                r_c2 = p_c.add_run(); r_c2.text = q_data["customerValue"]; r_c2.font.color.rgb = COLOR_TEXT_PRIMARY

                p_b = tf.add_paragraph()
                p_b.font.size = FONT_BODY; p_b.line_spacing = 1.14
                r_b1 = p_b.add_run(); r_b1.text = "Business Value: "; r_b1.font.bold = True; r_b1.font.color.rgb = COLOR_MYNTRA_GREEN
                r_b2 = p_b.add_run(); r_b2.text = q_data["businessValue"]; r_b2.font.color.rgb = COLOR_TEXT_PRIMARY
            else:
                p_ans = tf.add_paragraph()
                p_ans.text = q_data["ans"]
                p_ans.font.size = FONT_BODY; p_ans.line_spacing = 1.14; p_ans.font.color.rgb = COLOR_TEXT_PRIMARY

    # ==========================================
    # SLIDE 7: SOLUTION PRINCIPLES & STRUCTURED RICE TABLE
    # ==========================================
    elif s_num == 7:
        # Top Card: 3 Core Product Principles
        p_top = card_top
        p_h = Inches(1.4)
        add_card_shape(slide, Inches(0.65), p_top, Inches(12.033), p_h, COLOR_SYNTHESIS_BG, COLOR_SYNTHESIS_BORDER)
        box_p = slide.shapes.add_textbox(Inches(0.77), p_top + Inches(0.06), Inches(11.793), p_h - Inches(0.12))
        tf_p = box_p.text_frame; tf_p.word_wrap = True
        tf_p.margin_left = Inches(0.06); tf_p.margin_right = Inches(0.06)
        p = tf_p.paragraphs[0]
        p.text = "✨ 3 CORE PRODUCT DESIGN PRINCIPLES"
        p.font.size = FONT_CARD_HEADER; p.font.bold = True; p.font.color.rgb = COLOR_PURPLE_PILL; p.space_after = Pt(3)
        for princ_title, princ_desc in data["principles"]:
            p_item = tf_p.add_paragraph()
            p_item.font.size = FONT_BODY; p_item.line_spacing = 1.12
            p_item.space_before = Pt(1); p_item.space_after = Pt(2)
            set_bullet_indent(p_item, left_margin_pt=12, hanging_indent_pt=0)
            r1 = p_item.add_run(); r1.text = "★ " + princ_title + ": "; r1.font.bold = True; r1.font.color.rgb = COLOR_MYNTRA_PINK
            r2 = p_item.add_run(); r2.text = princ_desc; r2.font.color.rgb = COLOR_TEXT_PRIMARY

        # Bottom: Structured RICE Prioritization Table
        t_top = p_top + p_h + Inches(0.14)
        t_h = card_height - p_h - Inches(0.14)
        add_card_shape(slide, Inches(0.65), t_top, Inches(12.033), t_h, COLOR_BG_CARD, COLOR_BORDER)
        box_t = slide.shapes.add_textbox(Inches(0.77), t_top + Inches(0.08), Inches(11.793), t_h - Inches(0.16))
        tf_t = box_t.text_frame; tf_t.word_wrap = True
        tf_t.margin_left = Inches(0.06); tf_t.margin_right = Inches(0.06)
        
        p = tf_t.paragraphs[0]
        p.text = "📊 QUANTITATIVE RICE PRIORITIZATION MATRIX & ROADMAP"
        p.font.size = FONT_CARD_HEADER; p.font.bold = True; p.font.color.rgb = COLOR_PURPLE_PILL; p.space_after = Pt(4)

        for row in data["riceTable"]:
            p_row = tf_t.add_paragraph()
            p_row.font.size = FONT_BODY; p_row.line_spacing = 1.15
            p_row.space_before = Pt(2); p_row.space_after = Pt(3)
            set_bullet_indent(p_row, left_margin_pt=14, hanging_indent_pt=0)
            
            # Solution Pillar
            r_sol = p_row.add_run(); r_sol.text = row["solution"] + ": "; r_sol.font.bold = True; r_sol.font.color.rgb = COLOR_MYNTRA_PINK if "Winner" in row["verdict"] else COLOR_TEXT_PRIMARY
            r_desc = p_row.add_run(); r_desc.text = row["desc"] + " \u2014 "; r_desc.font.color.rgb = COLOR_TEXT_PRIMARY
            
            # RICE Math
            r_math = p_row.add_run()
            r_math.text = f"Reach: {row['reach']} | Impact: {row['impact']} | Conf: {row['confidence']} | Effort: {row['effort']} \u2192 "
            r_math.font.color.rgb = COLOR_TEXT_MUTED
            
            # Score & Status
            r_score = p_row.add_run(); r_score.text = f"Score: {row['score']} "; r_score.font.bold = True; r_score.font.color.rgb = COLOR_MYNTRA_PINK
            r_stat = p_row.add_run(); r_stat.text = f"[{row['verdict']}]"; r_stat.font.bold = True; r_stat.font.color.rgb = COLOR_MYNTRA_GREEN if "Winner" in row["verdict"] else COLOR_TEXT_MUTED

    # ==========================================
    # SLIDE 8: ARCHITECTURE & MVP FEATURE DEMO
    # ==========================================
    elif s_num == 8:
        # Top: 4-Step Architecture Pipeline
        pipe_h = Inches(1.6)
        add_card_shape(slide, Inches(0.65), card_top, Inches(12.033), pipe_h, COLOR_BG_CARD, COLOR_BORDER)
        box_p = slide.shapes.add_textbox(Inches(0.77), card_top + Inches(0.06), Inches(11.793), pipe_h - Inches(0.12))
        tf_p = box_p.text_frame; tf_p.word_wrap = True
        tf_p.margin_left = Inches(0.06); tf_p.margin_right = Inches(0.06)
        p = tf_p.paragraphs[0]
        p.text = "🏗️ TECHNICAL ARCHITECTURE PIPELINE (LATENCY <180ms)"
        p.font.size = FONT_CARD_HEADER; p.font.bold = True; p.font.color.rgb = COLOR_PURPLE_PILL; p.space_after = Pt(2)

        p_arch = tf_p.add_paragraph()
        p_arch.font.size = FONT_BODY; p_arch.line_spacing = 1.15
        for s_idx, (s_title, s_desc) in enumerate(data["pipeline"]):
            r0 = p_arch.add_run(); r0.text = f"[{s_title}]: "; r0.font.bold = True; r0.font.color.rgb = COLOR_MYNTRA_PINK
            r1 = p_arch.add_run(); r1.text = s_desc + (" \u2192 " if s_idx < len(data["pipeline"])-1 else "")
            r1.font.color.rgb = COLOR_TEXT_PRIMARY

        # Bottom: 3 Live MVP Features
        mvp_top = card_top + pipe_h + Inches(0.14)
        mvp_w = Inches(3.9)
        mvp_gap = Inches(0.16)
        mvp_h = card_height - pipe_h - Inches(0.14)
        for m_idx, feat in enumerate(data["mvpFeatures"]):
            m_left = Inches(0.65) + m_idx * (mvp_w + mvp_gap)
            add_card_shape(slide, m_left, mvp_top, mvp_w, mvp_h, COLOR_SYNTHESIS_BG, COLOR_SYNTHESIS_BORDER)
            box_m = slide.shapes.add_textbox(m_left + Inches(0.08), mvp_top + Inches(0.06), mvp_w - Inches(0.16), mvp_h - Inches(0.12))
            tf_m = box_m.text_frame; tf_m.word_wrap = True
            tf_m.margin_left = Inches(0.06); tf_m.margin_right = Inches(0.06)
            
            p = tf_m.paragraphs[0]
            p.text = "📱 " + feat["title"]
            p.font.size = FONT_CARD_HEADER; p.font.bold = True; p.font.color.rgb = COLOR_PURPLE_PILL; p.space_after = Pt(2)
            
            p_cap = tf_m.add_paragraph()
            p_cap.font.size = FONT_BODY; p_cap.line_spacing = 1.12; p_cap.space_after = Pt(2)
            r1 = p_cap.add_run(); r1.text = feat["desc"]; r1.font.color.rgb = COLOR_TEXT_PRIMARY
            
            # Embedded Screenshot if exists
            img_p = feat.get("figmaImage")
            if img_p and os.path.exists(img_p):
                slide.shapes.add_picture(img_p, m_left + Inches(0.2), mvp_top + Inches(1.15), mvp_w - Inches(0.4), mvp_h - Inches(1.25))

    # ==========================================
    # SLIDE 9: METRICS HIERARCHY & EXPERIMENT DESIGN
    # ==========================================
    elif s_num == 9:
        # Left: Structured Success Metrics Hierarchy Table
        left_w = Inches(6.8)
        add_card_shape(slide, Inches(0.65), card_top, left_w, card_height, COLOR_BG_CARD, COLOR_BORDER)
        box_m = slide.shapes.add_textbox(Inches(0.77), card_top + Inches(0.08), left_w - Inches(0.24), card_height - Inches(0.16))
        tf_m = box_m.text_frame; tf_m.word_wrap = True
        tf_m.margin_left = Inches(0.06); tf_m.margin_right = Inches(0.06)
        
        p = tf_m.paragraphs[0]
        p.text = "🎯 SUCCESS METRICS HIERARCHY"
        p.font.size = FONT_CARD_HEADER; p.font.bold = True; p.font.color.rgb = COLOR_PURPLE_PILL; p.space_after = Pt(4)

        for row in data["metricsTable"]:
            p_row = tf_m.add_paragraph()
            p_row.font.size = FONT_BODY; p_row.line_spacing = 1.15
            p_row.space_before = Pt(2); p_row.space_after = Pt(3)
            set_bullet_indent(p_row, left_margin_pt=14, hanging_indent_pt=0)
            
            # Type & Metric Name
            r_type = p_row.add_run(); r_type.text = "[" + row["type"] + "] "; r_type.font.bold = True; r_type.font.color.rgb = COLOR_MYNTRA_PINK if "North Star" in row["type"] else COLOR_PURPLE_PILL
            r_name = p_row.add_run(); r_name.text = row["kpi"] + ": "; r_name.font.bold = True; r_name.font.color.rgb = COLOR_TEXT_PRIMARY
            
            # Baseline with Citation
            r_tgt = p_row.add_run(); r_tgt.text = row["target"] + " "; r_tgt.font.bold = True; r_tgt.font.color.rgb = COLOR_MYNTRA_GREEN
            
            # Guardrail / Goal
            r_goal = p_row.add_run(); r_goal.text = f"({row['goal']})"; r_goal.font.color.rgb = COLOR_TEXT_MUTED

        # Right: 200,000-User RCT A/B Testing Design & Schema
        right_left = Inches(7.6)
        right_w = Inches(5.083)
        add_card_shape(slide, right_left, card_top, right_w, card_height, COLOR_SYNTHESIS_BG, COLOR_SYNTHESIS_BORDER)
        box_exp = slide.shapes.add_textbox(right_left + Inches(0.1), card_top + Inches(0.08), right_w - Inches(0.2), card_height - Inches(0.16))
        tf_exp = box_exp.text_frame; tf_exp.word_wrap = True
        tf_exp.margin_left = Inches(0.06); tf_exp.margin_right = Inches(0.06)
        
        p = tf_exp.paragraphs[0]
        p.text = "🧪 200,000-USER RCT A/B TESTING DESIGN"
        p.font.size = FONT_CARD_HEADER; p.font.bold = True; p.font.color.rgb = COLOR_PURPLE_PILL; p.space_after = Pt(4)

        for b_lbl, b_txt in data["experimentDesign"]:
            p_item = tf_exp.add_paragraph()
            p_item.font.size = FONT_BODY; p_item.line_spacing = 1.15
            p_item.space_before = Pt(2); p_item.space_after = Pt(3)
            set_bullet_indent(p_item, left_margin_pt=14, hanging_indent_pt=10)
            r1 = p_item.add_run(); r1.text = "• " + b_lbl + " "; r1.font.bold = True; r1.font.color.rgb = COLOR_TEXT_PRIMARY
            r2 = p_item.add_run(); r2.text = b_txt; r2.font.color.rgb = COLOR_TEXT_PRIMARY

    # ==========================================
    # SLIDE 10: RISK MATRIX & 3-PHASE GTM ROLLOUT
    # ==========================================
    elif s_num == 10:
        # Top: 3 Pitfall & Mitigation Risk Cards
        pit_w = Inches(3.9)
        pit_gap = Inches(0.16)
        pit_h = (card_height - Inches(0.18)) / 2.0
        for p_idx, pit in enumerate(data["pitfalls"]):
            p_left = Inches(0.65) + p_idx * (pit_w + pit_gap)
            add_card_shape(slide, p_left, card_top, pit_w, pit_h, COLOR_BG_CARD, COLOR_BORDER)
            box = slide.shapes.add_textbox(p_left + Inches(0.08), card_top + Inches(0.06), pit_w - Inches(0.16), pit_h - Inches(0.12))
            tf = box.text_frame; tf.word_wrap = True
            tf.margin_left = Inches(0.06); tf.margin_right = Inches(0.06)
            
            p = tf.paragraphs[0]
            p.text = "⚠️ " + pit["title"]
            p.font.size = FONT_CARD_HEADER; p.font.bold = True; p.font.color.rgb = COLOR_MYNTRA_PINK; p.space_after = Pt(2)
            
            p_imp = tf.add_paragraph()
            p_imp.font.size = FONT_BODY; p_imp.line_spacing = 1.14; p_imp.space_after = Pt(2)
            set_bullet_indent(p_imp, left_margin_pt=10, hanging_indent_pt=0)
            r1 = p_imp.add_run(); r1.text = "Pitfall: "; r1.font.bold = True; r1.font.color.rgb = COLOR_TEXT_PRIMARY
            r2 = p_imp.add_run(); r2.text = pit["pitfall"]; r2.font.color.rgb = COLOR_TEXT_PRIMARY
            
            p_mit = tf.add_paragraph()
            p_mit.font.size = FONT_BODY; p_mit.line_spacing = 1.14; p_mit.space_after = Pt(0)
            set_bullet_indent(p_mit, left_margin_pt=10, hanging_indent_pt=0)
            r3 = p_mit.add_run(); r3.text = "Mitigation: "; r3.font.bold = True; r3.font.color.rgb = COLOR_MYNTRA_GREEN
            r4 = p_mit.add_run(); r4.text = pit["mitigation"]; r4.font.color.rgb = COLOR_TEXT_PRIMARY

        # Bottom: 3-Phase Rollout Roadmap
        r_top = card_top + pit_h + Inches(0.18)
        r_card_w = Inches(3.9)
        r_card_gap = Inches(0.16)
        r_card_h = (card_height - Inches(0.18)) / 2.0
        for r_idx, phase in enumerate(data["rolloutPhases"]):
            r_left = Inches(0.65) + r_idx * (r_card_w + r_card_gap)
            add_card_shape(slide, r_left, r_top, r_card_w, r_card_h, COLOR_SYNTHESIS_BG, COLOR_SYNTHESIS_BORDER)
            box = slide.shapes.add_textbox(r_left + Inches(0.08), r_top + Inches(0.06), r_card_w - Inches(0.16), r_card_h - Inches(0.12))
            tf = box.text_frame; tf.word_wrap = True
            tf.margin_left = Inches(0.06); tf.margin_right = Inches(0.06)
            
            p = tf.paragraphs[0]
            p.text = "🚀 " + phase["phase"]
            p.font.size = FONT_CARD_HEADER; p.font.bold = True; p.font.color.rgb = COLOR_PURPLE_PILL; p.space_after = Pt(2)
            
            p_tgt = tf.add_paragraph()
            p_tgt.font.size = FONT_BODY; p_tgt.line_spacing = 1.14; p_tgt.space_after = Pt(2)
            set_bullet_indent(p_tgt, left_margin_pt=10, hanging_indent_pt=0)
            r1 = p_tgt.add_run(); r1.text = "Target: "; r1.font.bold = True; r1.font.color.rgb = COLOR_TEXT_PRIMARY
            r2 = p_tgt.add_run(); r2.text = phase["target"]; r2.font.color.rgb = COLOR_TEXT_PRIMARY
            
            p_scp = tf.add_paragraph()
            p_scp.font.size = FONT_BODY; p_scp.line_spacing = 1.14; p_scp.space_after = Pt(2)
            set_bullet_indent(p_scp, left_margin_pt=10, hanging_indent_pt=0)
            r3 = p_scp.add_run(); r3.text = "Scope: "; r3.font.bold = True; r3.font.color.rgb = COLOR_TEXT_PRIMARY
            r4 = p_scp.add_run(); r4.text = phase["scope"]; r4.font.color.rgb = COLOR_TEXT_PRIMARY
            
            p_gate = tf.add_paragraph()
            p_gate.font.size = FONT_BODY; p_gate.line_spacing = 1.14; p_gate.space_after = Pt(0)
            set_bullet_indent(p_gate, left_margin_pt=10, hanging_indent_pt=0)
            r5 = p_gate.add_run(); r5.text = "Gate: "; r5.font.bold = True; r5.font.color.rgb = COLOR_MYNTRA_GREEN
            r6 = p_gate.add_run(); r6.text = phase["gate"]; r6.font.color.rgb = COLOR_TEXT_PRIMARY

# Save Presentation to both filenames
output_path_primary = "Myntra_Wishlist_Studio_10_Slide_Deck.pptx"
output_path_updated = "Myntra_Wishlist_Studio_10_Slide_Deck_Updated.pptx"

try:
    prs.save(output_path_primary)
    print(f"SUCCESS: Generated 16:9 Executive PowerPoint Presentation at {output_path_primary}")
except PermissionError:
    print(f"NOTICE: {output_path_primary} is locked.")

try:
    prs.save(output_path_updated)
    print(f"SUCCESS: Saved updated presentation at {output_path_updated}")
except PermissionError:
    print(f"NOTICE: {output_path_updated} is locked.")
